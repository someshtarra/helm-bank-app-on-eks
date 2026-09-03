#!/usr/bin/env python3
"""
Generate an executive-ready, 21-slide PowerPoint presentation (.pptx)
for the Bank Application on AWS EKS DevOps portfolio project.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

os.makedirs('presentation', exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# Theme Colors
C_DARK = RGBColor(15, 23, 42)       # Slate 900
C_NAVY = RGBColor(30, 41, 59)       # Slate 800
C_BLUE = RGBColor(50, 108, 229)     # Kubernetes Blue
C_AWS = RGBColor(255, 153, 0)       # AWS Orange
C_TEAL = RGBColor(0, 137, 123)      # Teal / Security
C_PURPLE = RGBColor(126, 87, 194)   # Purple / Metrics
C_STORAGE = RGBColor(46, 125, 50)   # Storage Green
C_BG_LIGHT = RGBColor(248, 250, 252)# Light Gray / Slate 50
C_WHITE = RGBColor(255, 255, 255)
C_MUTED = RGBColor(100, 116, 139)   # Slate 500
C_CARD_BORDER = RGBColor(203, 213, 225)

def add_header(slide, title_text, category_text="BANK APPLICATION ON AWS EKS — DEVOPS ARCHITECTURE"):
    # Header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_DARK
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = Inches(0.8)
    tf.margin_top = Inches(0.18)

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(9.5)
    p_cat.font.bold = True
    p_cat.font.color.rgb = C_AWS

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = C_WHITE

def add_card(slide, left, top, width, height, title, content_bullets, border_color=C_CARD_BORDER, bg_color=C_WHITE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)

    tf = card.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.25)
    tf.word_wrap = True

    p_title = tf.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(14)
    p_title.font.bold = True
    p_title.font.color.rgb = C_DARK
    p_title.space_after = Pt(10)

    for bullet in content_bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(11)
        p.font.color.rgb = C_NAVY
        p.space_after = Pt(6)

def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = notes_text

# ------------------------------------------------------------------------------
# Slide 1: Title Slide
# ------------------------------------------------------------------------------
slide1 = prs.slides.add_slide(blank_layout)
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = C_DARK
bg1.line.fill.background()

title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
tf1 = title_box.text_frame
p_sub = tf1.paragraphs[0]
p_sub.text = "ENTERPRISE CLOUD-NATIVE WORKLOAD ARCHITECTURE"
p_sub.font.size = Pt(14)
p_sub.font.bold = True
p_sub.font.color.rgb = C_AWS
p_sub.space_after = Pt(14)

p_main = tf1.add_paragraph()
p_main.text = "Bank Application on Kubernetes (AWS EKS)"
p_main.font.size = Pt(36)
p_main.font.bold = True
p_main.font.color.rgb = C_WHITE
p_main.space_after = Pt(14)

p_desc = tf1.add_paragraph()
p_desc.text = "Production Readiness Review, Horizontal Elasticity (HPA v2), Zero-Trust RBAC, and Resilient Storage Architecture"
p_desc.font.size = Pt(16)
p_desc.font.color.rgb = C_MUTED
p_desc.space_after = Pt(24)

p_meta = tf1.add_paragraph()
p_meta.text = "Author: Senior DevOps Engineer & Kubernetes Instructor  |  Helm v3 & AWS EKS"
p_meta.font.size = Pt(12)
p_meta.font.color.rgb = C_BLUE

add_speaker_notes(slide1, "Welcome everyone. Today we are walking through the complete production transformation of the Bank Application on AWS EKS. We will explore architectural flaws in traditional setups, how we fixed them, and the production-grade patterns implemented across autoscaling, storage, security, and lifecycle management.")

# ------------------------------------------------------------------------------
# Slide 2: Business & Technical Problem
# ------------------------------------------------------------------------------
slide2 = prs.slides.add_slide(blank_layout)
add_header(slide2, "1. Business & Technical Problem Statement")
add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "The Business Challenge",
         [
             "High Transaction Volatility: Traffic surges during market openings and paydays demand instant elasticity.",
             "Strict Zero-Downtime SLA: In financial banking, service disruptions cause immediate financial and brand loss.",
             "Regulatory Compliance: Strict mandates (PCI-DSS, SOC2) require zero-trust least-privilege access and isolation.",
             "Cost Governance: Over-provisioning static EC2 instances drives up AWS cloud spend needlessly."
         ], C_AWS)

add_card(slide2, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "The Technical Challenge",
         [
             "Legacy YAML Sprawl: Unparameterized, static manifests hardcoded to wrong namespaces ('default' vs 'lab').",
             "Storage Contention: Multiple replicas mounted to a single AWS EBS volume (ReadWriteOnce) causing Multi-Attach failures.",
             "Autoscaling Breakdown: HPA v1 misconfigured without mandatory CPU requests, preventing dynamic scaling.",
             "Security Exposure: Containers running as root with excessive cluster pod deletion privileges."
         ], C_BLUE)
add_speaker_notes(slide2, "Every architectural design must solve real business problems. Banking systems face intense spikes and strict compliance. If infrastructure is poorly structured, outages occur during scaling and security vulnerabilities emerge.")

# ------------------------------------------------------------------------------
# Slide 3: Technology Stack
# ------------------------------------------------------------------------------
slide3 = prs.slides.add_slide(blank_layout)
add_header(slide3, "2. Technology Stack & Enterprise Tooling")
techs = [
    ("AWS EKS", "Managed Kubernetes Control Plane with 99.95% SLA, multi-AZ etcd replication, and VPC CNI integration.", C_AWS),
    ("Helm v3", "Kubernetes package management, parameterization via values.yaml, and release versioning.", C_BLUE),
    ("Kubernetes HPA v2", "Dynamic horizontal elasticity with dual CPU/Memory metrics and stabilization behavior.", C_PURPLE),
    ("AWS EBS CSI Driver", "Dynamic gp3 block storage provisioning with WaitForFirstConsumer AZ binding.", C_TEAL),
    ("AWS Network LB", "Ultra-low latency Layer 4 load balancing with cross-zone failover and instance targets.", C_AWS),
    ("Apache HTTP Server", "Hardened web server container running as unprivileged UID 10001 with dropped capabilities.", C_NAVY)
]
for i, (name, desc, col) in enumerate(techs):
    x = Inches(0.8 + (i % 3) * 3.9)
    y = Inches(1.5 + (i // 3) * 2.7)
    add_card(slide3, x, y, Inches(3.6), Inches(2.4), name, [desc], col)
add_speaker_notes(slide3, "Our technology stack combines battle-tested cloud-native components: AWS EKS for control plane reliability, Helm v3 for configuration as code, HPA v2 for elasticity, and modern EBS CSI for persistent storage.")

# ------------------------------------------------------------------------------
# Slide 4: High-Level Architecture (Diagram)
# ------------------------------------------------------------------------------
slide4 = prs.slides.add_slide(blank_layout)
add_header(slide4, "3. High-Level Architecture & End-to-End Traffic Flow")
if os.path.exists('diagrams/architecture.png'):
    slide4.shapes.add_picture('diagrams/architecture.png', Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.6))
add_card(slide4, Inches(9.6), Inches(1.4), Inches(3.0), Inches(5.6),
         "Core Flow Summary",
         [
             "Traffic Ingress: User -> AWS NLB -> K8s Service -> Pods.",
             "Scheduling: Dedicated frontend nodes via taints & tolerations.",
             "Elasticity: HPA v2 queries Metrics Server to autoscale pods 2 to 10.",
             "Storage: Dynamic EBS gp3 volume attached via AWS CSI driver.",
             "RBAC: Least-privilege ServiceAccount with token automount disabled."
         ], C_BLUE)
add_speaker_notes(slide4, "Here is the comprehensive architectural diagram. Notice how client requests flow through the AWS Network Load Balancer directly into the Kubernetes Service, which round-robins across healthy pods. Concurrently, HPA monitors workload utilization.")

# ------------------------------------------------------------------------------
# Slide 5: AWS EKS Infrastructure Topology
# ------------------------------------------------------------------------------
slide5 = prs.slides.add_slide(blank_layout)
add_header(slide5, "4. AWS EKS Cluster Topology & Isolation")
add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Node Group Architecture",
         [
             "Multi-AZ Deployment: Worker nodes distributed across us-east-1a and us-east-1b for high availability.",
             "Workload Taints: Dedicated node groups tainted with 'dedicated=frontend:NoSchedule'.",
             "Node Affinity: Pods declare soft affinity ('preferredDuringScheduling') to production-labeled instances.",
             "Auto-Healing: AWS EKS Managed Node Groups automatically replace degraded EC2 instances."
         ], C_AWS)
add_card(slide5, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Networking & Security Perimeter",
         [
             "AWS VPC CNI: Native VPC IP allocation per pod, eliminating overlay network translation overhead.",
             "Security Groups for Pods: Direct ENI binding allowing granular firewall rules per microservice.",
             "Private API Server: EKS control plane accessible via VPN/DirectConnect with public endpoint restricted.",
             "AWS OIDC Identity: Native federation enabling IAM Roles for Service Accounts (IRSA)."
         ], C_TEAL)
add_speaker_notes(slide5, "In EKS, nodes are separated by Availability Zones. We isolate frontend banking pods on dedicated worker nodes using taints and tolerations to protect them from resource contention.")

# ------------------------------------------------------------------------------
# Slide 6: Helm Chart Structure
# ------------------------------------------------------------------------------
slide6 = prs.slides.add_slide(blank_layout)
add_header(slide6, "5. Helm Chart Structure & Templating Best Practices")
add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Standardized Repository Hierarchy",
         [
             "Chart.yaml: Semantic chart metadata (v2 API, appVersion 2.4.58).",
             "values.yaml: Self-documenting, production-grade default values.",
             "templates/_helpers.tpl: Reusable template functions for labels and naming.",
             "templates/*.yaml: Declarative manifests for Deployment, Service, HPA, RBAC, Storage, and PDB.",
             "NOTES.txt: Post-installation verification commands."
         ], C_BLUE)
add_card(slide6, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Templating Improvements",
         [
             "Dynamic Namespaces: Eliminated hardcoded 'default' and 'lab'; standardized on {{ .Release.Namespace }}.",
             "Standard Labeling: Applied Kubernetes recommended labels ('app.kubernetes.io/*') across all objects.",
             "Schema Compliance: Fixed roleRef schema from invalid list to correct map syntax.",
             "Modular Toggles: Configurable flags for storage, autoscaling, RBAC, and security contexts."
         ], C_DARK)
add_speaker_notes(slide6, "We refactored the Helm chart into an enterprise layout with _helpers.tpl for centralized naming and standardized labels. Hardcoded namespaces were completely eliminated.")

# ------------------------------------------------------------------------------
# Slide 7: Deployment Architecture
# ------------------------------------------------------------------------------
slide7 = prs.slides.add_slide(blank_layout)
add_header(slide7, "6. Deployment Architecture & Pod Lifecycle")
add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Deployment Specification",
         [
             "Target Image: someshtarra/httpd:fox-3D (Apache HTTP Server).",
             "Resource QoS: Configured with Burstable QoS class (200m CPU / 128Mi RAM requests).",
             "Non-Root User: Runs explicitly as UID 10001 with GID 10001.",
             "RollingUpdate: maxSurge: 1, maxUnavailable: 0 to guarantee zero downtime during releases."
         ], C_BLUE)
add_card(slide7, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Pod Resilience & Placement",
         [
             "Tolerations: Matches 'dedicated=frontend:NoSchedule' taints.",
             "PodDisruptionBudget: Guarantees at least 1 pod is always available during node maintenance.",
             "Volume Mounts: Conditionally mounts persistent storage at /var/www/html/.",
             "Topology Spread: Distributes pods across distinct nodes and availability zones."
         ], C_TEAL)
add_speaker_notes(slide7, "The deployment manifest controls the pod lifecycle. We enforce non-root execution, specify precise compute requests, and integrate a PodDisruptionBudget to protect pods during node draining.")

# ------------------------------------------------------------------------------
# Slide 8: Kubernetes Service & AWS Load Balancer
# ------------------------------------------------------------------------------
slide8 = prs.slides.add_slide(blank_layout)
add_header(slide8, "7. Kubernetes Service & AWS Load Balancer")
add_card(slide8, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Service Layer Configuration",
         [
             "Service Type: LoadBalancer (Port 80 -> TargetPort 80).",
             "Selector Match: Strictly aligned with deployment pod template labels.",
             "Endpoint Management: Endpoints automatically updated based on readinessProbe status.",
             "Multi-Port Support: Easily extendable to HTTPS (Port 443) with ACM certificate binding."
         ], C_BLUE)
add_card(slide8, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "AWS NLB Controller Annotations",
         [
             "aws-load-balancer-type: 'external' — Provisions modern AWS Network Load Balancer.",
             "aws-load-balancer-nlb-target-type: 'instance' — Routes directly to node instance targets.",
             "aws-load-balancer-scheme: 'internet-facing' — Enables public client connectivity.",
             "Performance Advantage: Layer 4 NLB handles millions of requests/sec with ultra-low latency."
         ], C_AWS)
add_speaker_notes(slide8, "By using modern AWS Load Balancer Controller annotations, we provision a high-performance Network Load Balancer rather than the legacy Classic Load Balancer.")

# ------------------------------------------------------------------------------
# Slide 9: HPA & Auto Scaling (Diagram)
# ------------------------------------------------------------------------------
slide9 = prs.slides.add_slide(blank_layout)
add_header(slide9, "8. Horizontal Pod Autoscaler (HPA v2) Mechanics")
if os.path.exists('diagrams/hpa-flow.png'):
    slide9.shapes.add_picture('diagrams/hpa-flow.png', Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.6))
add_card(slide9, Inches(9.6), Inches(1.4), Inches(3.0), Inches(5.6),
         "HPA Scaling Logic",
         [
             "Telemetry: Metrics Server scrapes kubelet /metrics/resource.",
             "Sync Period: Controller evaluates every 15 seconds.",
             "Target: Scales up when average CPU > 50% or RAM > 80%.",
             "Scale Range: Dynamically scales between 2 and 10 replicas.",
             "Stabilization: 300s scale-down delay prevents thrashing."
         ], C_PURPLE)
add_speaker_notes(slide9, "This diagram explains the complete HPA loop. The Metrics Server aggregates CPU usage from kubelets, and the HPA controller calculates target replicas and patches the deployment.")

# ------------------------------------------------------------------------------
# Slide 10: CPU Requests, Limits & Metrics Server Math
# ------------------------------------------------------------------------------
slide10 = prs.slides.add_slide(blank_layout)
add_header(slide10, "9. CPU Requests, Limits & Telemetry Mathematics")
add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "The HPA Formula & CPU Requests",
         [
             "Mathematical Formula: Desired = ceil(Current * (CurrentUsage / TargetUsage)).",
             "The Utilization Ratio: CPU % = (Actual Millicores / Requested Millicores) * 100.",
             "The Missing Request Pitfall: If CPU request is missing, denominator is undefined!",
             "Failure State: HPA displays '<unknown> / 50%' and cannot autoscale."
         ], C_PURPLE)
add_card(slide10, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Requests vs. Limits (QoS Mechanics)",
         [
             "Requests (200m CPU / 128Mi RAM): Guaranteed minimum compute reserved by scheduler.",
             "Limits (500m CPU / 256Mi RAM): Hard upper ceiling enforced by Linux kernel cgroups.",
             "CPU Exceeded: Results in CFS quota CPU throttling (latency rises; process is NOT killed).",
             "Memory Exceeded: Results in kernel OOM Killer terminating container (Exit Code 137)."
         ], C_BLUE)
add_speaker_notes(slide10, "Understanding the math behind HPA is vital. CPU percentage is calculated against the CPU request. Without requests, HPA fails completely. We also contrast CPU throttling with memory OOMKills.")

# ------------------------------------------------------------------------------
# Slide 11: Persistent Storage with AWS EBS (Diagram)
# ------------------------------------------------------------------------------
slide11 = prs.slides.add_slide(blank_layout)
add_header(slide11, "10. Storage Architecture: AWS EBS & CSI Driver")
if os.path.exists('diagrams/storage-flow.png'):
    slide11.shapes.add_picture('diagrams/storage-flow.png', Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.6))
add_card(slide11, Inches(9.6), Inches(1.4), Inches(3.0), Inches(5.6),
         "Storage Components",
         [
             "StorageClass: ebs.csi.aws.com with type: gp3.",
             "Binding Mode: WaitForFirstConsumer ensures volume matches node AZ.",
             "PVC: Requests 5Gi with ReadWriteOnce access mode.",
             "ReclaimPolicy: Retain preserves disk upon PVC deletion.",
             "Mount Path: /var/www/html/ for web assets."
         ], C_STORAGE)
add_speaker_notes(slide11, "Persistent storage is managed by the AWS EBS CSI driver. We use gp3 volumes and WaitForFirstConsumer binding to ensure the volume is created in the exact AZ where the pod lands.")

# ------------------------------------------------------------------------------
# Slide 12: EBS vs. EFS / RWO vs. RWX Architectural Analysis
# ------------------------------------------------------------------------------
slide12 = prs.slides.add_slide(blank_layout)
add_header(slide12, "11. Architectural Trade-Off: AWS EBS vs. AWS EFS")
add_card(slide12, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "AWS EBS: ReadWriteOnce (RWO)",
         [
             "Storage Type: Block Device (Virtual SAN hard drive).",
             "Node Attachment: Strictly limited to 1 EC2 instance at a time.",
             "Failure Mode: Replicas on Node B fail with 'Multi-Attach error'.",
             "Best Use Case: Single-replica workloads, databases (PostgreSQL, MySQL), or StatefulSets with dedicated volume per pod."
         ], C_AWS)
add_card(slide12, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "AWS EFS: ReadWriteMany (RWX)",
         [
             "Storage Type: Managed Distributed NFSv4.1 File System.",
             "Node Attachment: Concurrently mountable by 1000+ pods across all AZs.",
             "Horizontal Scaling: Seamlessly scales alongside HPA replicas (2 to 10).",
             "Best Use Case: Multi-replica web frontends, shared assets, CMS, and stateless cloud applications."
         ], C_STORAGE)
add_speaker_notes(slide12, "This slide covers the fundamental storage trade-off. EBS cannot scale horizontally across multiple nodes because of the ReadWriteOnce limitation. For multi-replica deployments, AWS EFS (RWX) or S3 object storage is the proper production architecture.")

# ------------------------------------------------------------------------------
# Slide 13: RBAC & ServiceAccount (Diagram)
# ------------------------------------------------------------------------------
slide13 = prs.slides.add_slide(blank_layout)
add_header(slide13, "12. RBAC & ServiceAccount Governance")
if os.path.exists('diagrams/rbac-flow.png'):
    slide13.shapes.add_picture('diagrams/rbac-flow.png', Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.6))
add_card(slide13, Inches(9.6), Inches(1.4), Inches(3.0), Inches(5.6),
         "RBAC Highlights",
         [
             "Identity: Dedicated ServiceAccount per application.",
             "Token Security: automountServiceAccountToken: false.",
             "Syntax Fix: roleRef written as map with apiGroup singular.",
             "Namespace Fix: Scoped dynamically to release namespace.",
             "Least Privilege: Read-only access to ConfigMaps only."
         ], C_TEAL)
add_speaker_notes(slide13, "RBAC controls who can do what in the cluster. We resolved schema bugs in roleRef and eliminated dangerous pod-deletion permissions, securing the application identity.")

# ------------------------------------------------------------------------------
# Slide 14: Security Best Practices
# ------------------------------------------------------------------------------
slide14 = prs.slides.add_slide(blank_layout)
add_header(slide14, "13. Workload Hardening & DevSecOps Standards")
add_card(slide14, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Container Security Context",
         [
             "runAsNonRoot: true — Enforces execution under unprivileged UID 10001.",
             "allowPrivilegeEscalation: false — Blocks setuid binary privilege inheritance.",
             "capabilities.drop: ['ALL'] — Strips all standard Linux kernel capabilities.",
             "seccompProfile: RuntimeDefault — Blocks unauthorized system calls at the kernel layer."
         ], C_TEAL)
add_card(slide14, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Enterprise Secrets & Identity",
         [
             "AWS IRSA: Replaces static credentials with temporary IAM OIDC role assumption.",
             "External Secrets Operator: Syncs database credentials directly from AWS Secrets Manager.",
             "NetworkPolicies: Restricts ingress/egress traffic to legitimate service paths.",
             "Image Scanning: Automated vulnerability scanning via Trivy in CI/CD pipeline."
         ], C_DARK)
add_speaker_notes(slide14, "Our security posture adheres to CIS Kubernetes benchmarks and the Restricted Pod Security Standard. We run unprivileged containers and eliminate static credentials via IRSA.")

# ------------------------------------------------------------------------------
# Slide 15: Health Probes
# ------------------------------------------------------------------------------
slide15 = prs.slides.add_slide(blank_layout)
add_header(slide15, "14. Health Probes: Startup, Readiness, & Liveness")
probes_info = [
    ("Startup Probe", "Purpose: Verifies process startup.", "Config: HTTP GET on / (failureThreshold: 20, period: 5s).\nAllows up to 100 seconds for slow initialization.\nCrucial: Disables liveness probe while active to prevent premature killing.", C_BLUE),
    ("Readiness Probe", "Purpose: Controls traffic routing.", "Config: HTTP GET on / (failureThreshold: 3, period: 5s).\nFailing probe removes pod IP from Service endpoints.\nCrucial: Container remains running; zero user traffic routed to unhealthy pod.", C_TEAL),
    ("Liveness Probe", "Purpose: Deadlock recovery.", "Config: TCP socket check on port 80 (period: 10s).\nFailing probe triggers kubelet to kill and restart container.\nCrucial: Avoid deep database checks to prevent cascading cluster restarts.", C_PURPLE)
]
for i, (name, purp, body, col) in enumerate(probes_info):
    x = Inches(0.8 + i * 3.9)
    add_card(slide15, x, Inches(1.5), Inches(3.6), Inches(5.4), name, [purp, body], col)
add_speaker_notes(slide15, "Health probes form a critical triad: Startup protects initialization, Readiness shields users from incomplete responses, and Liveness recovers wedged processes.")

# ------------------------------------------------------------------------------
# Slide 16: Rolling Updates & Zero Downtime
# ------------------------------------------------------------------------------
slide16 = prs.slides.add_slide(blank_layout)
add_header(slide16, "15. Zero-Downtime Rolling Updates")
add_card(slide16, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "RollingUpdate Mechanics",
         [
             "maxSurge: 1 — Creates 1 new pod before terminating any existing pod.",
             "maxUnavailable: 0 — Guarantees 100% capacity remains online throughout rollout.",
             "Readiness Gate: Old pod is only terminated after new pod passes readiness probe.",
             "Traffic Continuity: Clients experience zero connection drops or 502 errors."
         ], C_BLUE)
add_card(slide16, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "PodDisruptionBudget (PDB)",
         [
             "minAvailable: 1 — Prevents Kubernetes from terminating all pods during node draining.",
             "Node Upgrades: Cluster administrator upgrades EKS nodes safely without application downtime.",
             "Voluntary Disruptions: Protects against accidental manual evictions and autoscaler drains.",
             "High Availability: Works in tandem with HPA to ensure uninterrupted banking operations."
         ], C_DARK)
add_speaker_notes(slide16, "Zero-downtime updates require maxUnavailable: 0 and readiness probes. The PodDisruptionBudget guarantees that voluntary evictions during cluster upgrades do not cause downtime.")

# ------------------------------------------------------------------------------
# Slide 17: Deployment Workflow
# ------------------------------------------------------------------------------
slide17 = prs.slides.add_slide(blank_layout)
add_header(slide17, "16. Operational Deployment Workflow")
add_card(slide17, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Pre-Deployment & Installation",
         [
             "1. Lint Chart: 'helm lint .' to validate syntax and chart rules.",
             "2. Local Render: 'helm template bank-app .' to inspect generated YAML.",
             "3. Dry-Run: 'kubectl apply --dry-run=server' to test against EKS API.",
             "4. Release Install: 'helm upgrade --install bank-app . --namespace production'."
         ], C_BLUE)
add_card(slide17, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Upgrades & Instant Rollback",
         [
             "Track History: 'helm history bank-app -n production' displays all revisions.",
             "Rolling Upgrade: 'helm upgrade bank-app . --set image.tag=v2.0.0'.",
             "Instant Rollback: 'helm rollback bank-app 1 -n production' in <5 seconds.",
             "Status Verification: 'kubectl rollout status deployment/bank-app'."
         ], C_AWS)
add_speaker_notes(slide17, "We follow a disciplined deployment lifecycle: linting, dry-run validation, atomic Helm releases, revision history tracking, and single-command rollbacks.")

# ------------------------------------------------------------------------------
# Slide 18: Troubleshooting Playbook
# ------------------------------------------------------------------------------
slide18 = prs.slides.add_slide(blank_layout)
add_header(slide18, "17. Production Troubleshooting Runbook")
incidents = [
    ("Multi-Attach Error", "Symptom: Pod stuck in ContainerCreating.\nCause: Single EBS volume mounted across nodes.\nFix: Migrate to AWS EFS or co-locate with pod affinity.", C_AWS),
    ("HPA <unknown> / 50%", "Symptom: HPA target shows <unknown>.\nCause: Missing CPU requests or Metrics Server dead.\nFix: Define resources.requests.cpu; restart Metrics Server.", C_PURPLE),
    ("OOMKilled (Exit 137)", "Symptom: Container restarts with Exit Code 137.\nCause: Container exceeded resources.limits.memory.\nFix: Increase memory limit or optimize application memory.", C_BLUE),
    ("Empty Endpoints", "Symptom: Service returns 502 Bad Gateway.\nCause: Selector mismatch or failed readiness probe.\nFix: Align labels; inspect readinessProbe endpoint.", C_TEAL)
]
for i, (name, desc, col) in enumerate(incidents):
    x = Inches(0.8 + (i % 2) * 5.8)
    y = Inches(1.5 + (i // 2) * 2.7)
    add_card(slide18, x, y, Inches(5.5), Inches(2.4), name, [desc], col)
add_speaker_notes(slide18, "Our troubleshooting guide provides actionable playbooks for the four most common production incidents: EBS multi-attach conflicts, HPA unknown metrics, OOM kills, and empty service endpoints.")

# ------------------------------------------------------------------------------
# Slide 19: Production Readiness Review
# ------------------------------------------------------------------------------
slide19 = prs.slides.add_slide(blank_layout)
add_header(slide19, "18. Production Readiness Review Audit")
add_card(slide19, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Before (Original Flaws)",
         [
             "❌ Inconsistent namespaces ('default' vs 'lab') breaking RBAC.",
             "❌ Broken roleRef schema syntax failing Kubernetes admission.",
             "❌ Frontend web container granted full pod deletion permissions.",
             "❌ Deprecated autoscaling/v1 lacking stabilization controls.",
             "❌ EBS RWO volume causing multi-node mount deadlock."
         ], C_AWS)
add_card(slide19, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "After (Production-Grade)",
         [
             "✅ Consistent namespace scoping via {{ .Release.Namespace }}.",
             "✅ Valid schema-compliant roleRef map syntax.",
             "✅ Least-privilege RBAC with token automounting disabled.",
             "✅ Upgraded to autoscaling/v2 with dual CPU/memory metrics.",
             "✅ Documented EBS limitations and defined AWS EFS blueprint.",
             "✅ Zero-downtime rolling update with PodDisruptionBudget."
         ], C_TEAL)
add_speaker_notes(slide19, "This audit slide demonstrates the value delivered: moving from a fragile, insecure setup to an enterprise-grade, hardened production workload.")

# ------------------------------------------------------------------------------
# Slide 20: Future DevOps Roadmap
# ------------------------------------------------------------------------------
slide20 = prs.slides.add_slide(blank_layout)
add_header(slide20, "19. Future DevOps & GitOps Platform Roadmap")
add_card(slide20, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "CI/CD & GitOps Automation",
         [
             "GitHub Actions: Automated linting, Trivy vulnerability scanning, and Helm packaging.",
             "Argo CD: Declarative GitOps deployment with automated drift detection and self-healing.",
             "Terraform: Infrastructure as Code for provisioning EKS clusters, VPCs, and IAM roles.",
             "Amazon ECR: Private, signed container registries with immutable digest pinning."
         ], C_BLUE)
add_card(slide20, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "Observability & Platform Governance",
         [
             "Prometheus & Grafana: Real-time RED metrics, cluster telemetry, and custom HPA triggers.",
             "Fluent Bit & Loki: Centralized application log aggregation and search.",
             "Kyverno / OPA Gatekeeper: Automated policy enforcement for Pod Security Standards.",
             "Karpenter: Intelligent, high-speed node autoscaling replacing Cluster Autoscaler."
         ], C_PURPLE)
add_speaker_notes(slide20, "Looking ahead, this project serves as the foundation for an end-to-end GitOps platform powered by Argo CD, GitHub Actions, Prometheus, and Karpenter.")

# ------------------------------------------------------------------------------
# Slide 21: DevOps Learning Summary
# ------------------------------------------------------------------------------
slide21 = prs.slides.add_slide(blank_layout)
add_header(slide21, "20. DevOps Learning Summary & Key Takeaways")
add_card(slide21, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4),
         "Core Technical Competencies",
         [
             "Production Kubernetes Architecture: Deep mastery of Pods, Deployments, Services, and Ingress.",
             "Elastic Autoscaling: Mathematical mechanics of HPA v2 and Metrics Server integration.",
             "Cloud Storage Realities: Navigating block vs. network file storage trade-offs in AWS.",
             "Zero-Trust Security: Enforcing least-privilege RBAC and container security contexts."
         ], C_BLUE)
add_card(slide21, Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.4),
         "DevOps Mindset & Value",
         [
             "Understanding WHY: Prioritizing architectural rationale over rote YAML configuration.",
             "Reliability by Default: Designing for failure with probes, PDBs, and zero-downtime rollouts.",
             "Clear Technical Communication: Ability to teach and document complex cloud architectures.",
             "Production Readiness: Bridging the gap between lab exercises and enterprise realities."
         ], C_AWS)
add_speaker_notes(slide21, "In conclusion, this project exemplifies the senior DevOps mindset: not just writing YAML, but understanding why configurations work, anticipating production failures, and securing workloads by default.")

prs.save('presentation/devops-kubernetes-project.pptx')
print("Successfully created presentation/devops-kubernetes-project.pptx (21 slides)")
